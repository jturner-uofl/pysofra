"""PowerPoint (PPTX) rendering via ``python-pptx``.

Writes a single-slide ``.pptx`` containing one table. The caller never
manipulates python-pptx shapes directly — pass a :class:`SofraTable` and
a path and the renderer takes care of slide creation, table sizing,
font, header shading, and footnote textbox.

The renderer is gated on the optional ``python-pptx`` dependency:

.. code-block:: text

   pip install pysofra[pptx]

If python-pptx isn't installed, calling ``.to_pptx`` raises
``ImportError`` with installation guidance.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from ..core.schema import HeaderRow, Row, SpanningHeader
from ..core.table import SofraTable
from ..themes.registry import resolve_theme


@dataclass
class PptxRenderer:
    """Write a SofraTable to a single-slide ``.pptx`` file."""

    slide_title: str | None = None
    slide_width_in: float = 13.333
    slide_height_in: float = 7.5

    def write(self, table: SofraTable, path: Path) -> Path:
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError as e:  # pragma: no cover
            raise ImportError(
                "PPTX export requires python-pptx. Install with "
                "`pip install pysofra[pptx]` or `pip install python-pptx`."
            ) from e

        theme = resolve_theme(table.theme_name)
        d = theme.pptx
        font_name: str = d.get("font_name", "Calibri")
        font_size: int = int(d.get("font_size", 14))
        header_fill: str = d.get("header_fill", "F2F2F2")
        # Theme-configurable text colours; sensible high-contrast
        # defaults so the header doesn't render as white-on-light
        # under PowerPoint's default table style.
        header_text_color: str = d.get("header_text_color", "1A202C")
        body_text_color: str = d.get("body_text_color", "1A202C")

        ncols = _ncols(table)
        n_header_rows = (1 if table.spanning_headers else 0) + len(table.headers)
        n_body_rows = len(table.rows)
        n_total_rows = max(1, n_header_rows + n_body_rows)

        prs = Presentation()
        prs.slide_width = Inches(self.slide_width_in)
        prs.slide_height = Inches(self.slide_height_in)
        blank = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank)

        # Optional slide title.
        title_text = self.slide_title or table.caption
        if title_text:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                Inches(self.slide_width_in - 1), Inches(0.6),
            )
            tf = title_box.text_frame
            tf.text = title_text
            run = tf.paragraphs[0].runs[0]
            run.font.name = font_name
            run.font.size = Pt(font_size + 6)
            run.font.bold = True

        # Table sizing: fit within the slide with margins.
        left = Inches(0.5)
        top_in = 1.1 if title_text else 0.5
        width_in = self.slide_width_in - 1
        avail_height_in = self.slide_height_in - (1.6 if title_text else 1.0)

        # Reserve vertical space for the footnotes textbox *before*
        # sizing the table, so footnotes never get clipped off the
        # bottom of the slide. PowerPoint applies a paragraph line
        # spacing roughly 1.5x the font height for body text; we use
        # 1.8x as a conservative reservation so the textbox never
        # clips its own contents even under PowerPoint's body-style
        # space-before/after defaults.
        fn_font_size = max(8, font_size - 4)
        fn_line_height_in = (fn_font_size * 1.8) / 72.0
        n_footnote_lines = len(table.footnotes)
        footnote_height_in = (
            n_footnote_lines * fn_line_height_in + 0.25
            if n_footnote_lines else 0.0
        )
        avail_height_in = max(1.0, avail_height_in - footnote_height_in)

        plot = getattr(table, "inline_plot", None)
        plot_png: bytes | None = (
            plot.png_bytes
            if plot is not None and getattr(plot, "png_bytes", None)
            else None
        )

        if plot_png is not None:
            # Reserve ~40% of the available vertical space for the plot.
            assert plot is not None  # paired with plot_png check above
            plot_height_in = min(float(plot.height_in), avail_height_in * 0.42)
            table_height_in = avail_height_in - plot_height_in - 0.15
        else:
            plot_height_in = 0.0
            table_height_in = avail_height_in

        if plot_png is not None and table.inline_svg_position == "above":
            import io
            slide.shapes.add_picture(
                io.BytesIO(plot_png),
                left, Inches(top_in),
                width=Inches(width_in),
                height=Inches(plot_height_in),
            )
            table_top_in = top_in + plot_height_in + 0.15
        else:
            table_top_in = top_in

        shape = slide.shapes.add_table(
            n_total_rows, ncols,
            left, Inches(table_top_in),
            Inches(width_in), Inches(table_height_in),
        )
        word_table = shape.table

        # PowerPoint applies a default "Medium Style 2 - Accent 1"
        # table style that paints the first row's text white. We
        # control every cell's fill and font colour explicitly, so
        # disable the style-driven overrides to avoid white-on-light
        # header text and zebra-striped body rows.
        import contextlib
        for attr in ("first_row", "first_col", "horz_banding",
                     "vert_banding", "last_row", "last_col"):
            with contextlib.suppress(AttributeError, ValueError):  # pragma: no cover - defensive
                setattr(word_table, attr, False)

        if plot_png is not None and table.inline_svg_position == "below":
            import io
            slide.shapes.add_picture(
                io.BytesIO(plot_png),
                left, Inches(table_top_in + table_height_in + 0.15),
                width=Inches(width_in),
                height=Inches(plot_height_in),
            )

        # Keep older variable names used downstream.
        top = Inches(table_top_in)
        width = Inches(width_in)
        height = Inches(table_height_in)

        row_idx = 0
        if table.spanning_headers:
            _write_spanning_row(
                word_table, row_idx, table.spanning_headers, ncols,
                font_name=font_name, font_size=font_size,
                header_fill=header_fill,
                header_text_color=header_text_color,
                RGBColor=RGBColor, Pt=Pt, PP_ALIGN=PP_ALIGN,
            )
            row_idx += 1
        for hr in table.headers:
            _write_header_row(
                word_table, row_idx, hr, ncols,
                font_name=font_name, font_size=font_size,
                header_fill=header_fill,
                header_text_color=header_text_color,
                RGBColor=RGBColor, Pt=Pt, PP_ALIGN=PP_ALIGN,
            )
            row_idx += 1
        for body_row in table.rows:
            _write_body_row(
                word_table, row_idx, body_row, ncols,
                font_name=font_name, font_size=font_size,
                body_text_color=body_text_color, RGBColor=RGBColor,
                Pt=Pt, PP_ALIGN=PP_ALIGN,
            )
            row_idx += 1

        # Footnotes textbox below the table. We:
        #   (1) reserve enough room ahead of time via footnote_height_in;
        #   (2) explicitly pin single-line spacing + zero space-before/after
        #       per paragraph so PowerPoint's body-text defaults don't
        #       silently inflate the rendered height;
        #   (3) enable SHAPE_TO_FIT_TEXT so the box grows if the
        #       heuristic still underestimates (belt and braces).
        if table.footnotes:
            from pptx.enum.text import MSO_AUTO_SIZE
            fn_top = top + height + Inches(0.05)
            fn_box = slide.shapes.add_textbox(
                left, fn_top, width, Inches(footnote_height_in),
            )
            tf = fn_box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            for i, fn in enumerate(table.footnotes):
                if i == 0:
                    tf.text = fn
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                    para.text = fn
                para.line_spacing = 1.0
                para.space_before = Pt(0)
                para.space_after = Pt(0)
                for run in para.runs:
                    run.font.name = font_name
                    run.font.size = Pt(fn_font_size)
                    run.font.italic = True
                    if RGBColor is not None:
                        # ``RGBColor.from_string`` is python-pptx's own
                        # constructor; its stubs are not strict, so an
                        # explicit type-ignore documents that we accept
                        # the third-party untyped call here.
                        run.font.color.rgb = RGBColor.from_string(  # type: ignore[no-untyped-call]
                            body_text_color
                        )

        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path))
        # python-pptx stamps every ZIP entry with the current wall-clock,
        # which breaks cross-process byte-determinism. Rewrite with fixed
        # entry mtimes so identical input always yields identical bytes.
        from ._zip_determinism import make_zip_deterministic
        make_zip_deterministic(path)
        return path

    def render(self, table: SofraTable) -> str:  # pragma: no cover
        raise NotImplementedError("PptxRenderer writes to disk; use .write(table, path).")


def _ncols(table: SofraTable) -> int:
    if table.headers:
        return len(table.headers[0].cells)
    if table.rows:
        return len(table.rows[0].cells)
    return 1


def _set_cell(cell: Any, text: str, *, bold: bool, italic: bool,
              font_name: str, font_size: int, align: str | None,
              fill_hex: str | None, font_color_hex: str | None,
              RGBColor: Any, Pt: Any, PP_ALIGN: Any) -> None:
    cell.text = ""
    para = cell.text_frame.paragraphs[0]
    if align == "right":
        para.alignment = PP_ALIGN.RIGHT
    elif align == "center":
        para.alignment = PP_ALIGN.CENTER
    else:
        para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    # Explicit colour. python-pptx's default table style paints header
    # text white; without an explicit override the header cell renders
    # white-on-light-fill, which is unreadable. We always set a colour
    # so the theme's choice is what ships, not PowerPoint's default.
    if font_color_hex is not None and RGBColor is not None:
        run.font.color.rgb = RGBColor.from_string(font_color_hex)
    if fill_hex is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(fill_hex)


def _write_header_row(word_table: Any, idx: int, hr: HeaderRow, ncols: int, *,
                      font_name: str, font_size: int, header_fill: str,
                      header_text_color: str,
                      RGBColor: Any, Pt: Any, PP_ALIGN: Any) -> None:
    row = word_table.rows[idx]
    for j, c in enumerate(hr.cells[:ncols]):
        cell = row.cells[j]
        text = c.text.replace("\n", "  ")
        _set_cell(
            cell, text, bold=True, italic=False,
            font_name=font_name, font_size=font_size,
            align=c.align, fill_hex=header_fill,
            font_color_hex=header_text_color,
            RGBColor=RGBColor, Pt=Pt, PP_ALIGN=PP_ALIGN,
        )


def _write_body_row(word_table: Any, idx: int, r: Row, ncols: int, *,
                    font_name: str, font_size: int,
                    body_text_color: str, RGBColor: Any,
                    Pt: Any, PP_ALIGN: Any) -> None:
    row = word_table.rows[idx]
    for j, c in enumerate(r.cells[:ncols]):
        cell = row.cells[j]
        text = c.text
        if c.indent > 0:
            text = "    " * c.indent + text
        _set_cell(
            cell, text, bold=c.bold or r.is_group_header, italic=c.italic,
            font_name=font_name, font_size=font_size,
            align=c.align, fill_hex=None,
            font_color_hex=body_text_color,
            RGBColor=RGBColor, Pt=Pt, PP_ALIGN=PP_ALIGN,
        )


def _write_spanning_row(word_table: Any, idx: int, spans: tuple[SpanningHeader, ...],
                        ncols: int, *, font_name: str, font_size: int,
                        header_fill: str, header_text_color: str,
                        RGBColor: Any, Pt: Any,
                        PP_ALIGN: Any) -> None:
    row = word_table.rows[idx]
    for span in spans:
        anchor = row.cells[span.start]
        # python-pptx supports cell merging via _tc/_tcPr is involved; use the
        # public ``merge`` API where available.
        for j in range(span.start + 1, span.end + 1):
            anchor.merge(row.cells[j])
        _set_cell(
            anchor, span.label, bold=True, italic=False,
            font_name=font_name, font_size=font_size, align="center",
            font_color_hex=header_text_color,
            fill_hex=header_fill, RGBColor=RGBColor, Pt=Pt, PP_ALIGN=PP_ALIGN,
        )
