"""Regression tests for PPTX slide-overflow defects.

The original bug (reported in May 2026): rendering a Table 1 with four
footnotes pushed the footnote textbox past the slide bottom, so the
last footnote ("SMD = standardized mean difference") was clipped off.
The fix in ``render/pptx.py`` reserves footnote-shaped vertical space
*before* sizing the table.

These tests open the rendered ``.pptx`` and verify every shape
(title, table, footnotes textbox, embedded plot if any) fits within
the slide's geometric bounds.
"""

from __future__ import annotations

from pathlib import Path

import numpy as np
import pandas as pd
import pytest

import pysofra as ps

EMU_PER_INCH = 914400
MAX_BOTTOM_OVERSHOOT_IN = 0.05  # tolerate < 50 mil rounding


def _shapes_in_bounds(pptx_path: Path) -> list[tuple[str, float, float]]:
    """Return shapes whose bottom edge exceeds the slide height.

    Each entry is ``(name, bottom_in, slide_h_in)``.
    """
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    slide_h = prs.slide_height / EMU_PER_INCH
    bad: list[tuple[str, float, float]] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            top = (shape.top or 0) / EMU_PER_INCH
            h   = (shape.height or 0) / EMU_PER_INCH
            bottom = top + h
            if bottom - slide_h > MAX_BOTTOM_OVERSHOOT_IN:
                bad.append((shape.name, bottom, slide_h))
    return bad


class TestPptxHeaderContrast:
    """Without explicit font-colour overrides PowerPoint's default
    table style (Medium Style 2 - Accent 1) paints the first row's
    text white, which is unreadable on our light-grey header fill.
    The fix sets every header cell's font colour explicitly and
    disables the default-style row flags."""

    def test_header_cells_have_explicit_dark_text(self, tmp_path):
        pytest.importorskip("pptx")
        from pptx import Presentation
        rng = np.random.default_rng(0)
        df = pd.DataFrame({
            "arm": rng.choice(["A", "B"], 50),
            "x":   rng.normal(size=50),
        })
        t = ps.tbl_one(df, by="arm").add_p()
        out = tmp_path / "t.pptx"
        t.to_pptx(out)
        prs = Presentation(str(out))
        tbl = next(s.table for s in prs.slides[0].shapes if s.has_table)
        # Every header cell run must have an explicitly-set rgb colour
        # (otherwise PowerPoint applies the default style's white).
        for j, cell in enumerate(tbl.rows[0].cells):
            run = cell.text_frame.paragraphs[0].runs[0]
            rgb = run.font.color.rgb
            assert rgb is not None, f"header col {j}: no explicit colour"
            # Dark — luminance should be well below mid-grey.
            r, g, b = int(str(rgb)[0:2], 16), int(str(rgb)[2:4], 16), int(str(rgb)[4:6], 16)
            luma = 0.2126*r + 0.7152*g + 0.0722*b
            assert luma < 100, (
                f"header col {j}: text colour #{rgb} is too light "
                f"(luminance {luma:.0f}) against the header fill"
            )

    def test_default_table_style_flags_disabled(self, tmp_path):
        pytest.importorskip("pptx")
        from pptx import Presentation
        df = pd.DataFrame({"arm": ["A", "B"] * 5, "x": range(10)})
        t = ps.tbl_one(df, by="arm")
        out = tmp_path / "t.pptx"
        t.to_pptx(out)
        prs = Presentation(str(out))
        tbl = next(s.table for s in prs.slides[0].shapes if s.has_table)
        # All five auto-style flags should be off so theming flows
        # purely from our renderer.
        for attr in ("first_row", "first_col",
                     "horz_banding", "vert_banding"):
            assert getattr(tbl, attr) is False, (
                f"table.{attr} = True; PowerPoint will override our "
                f"styling for that band"
            )


class TestPptxFootnotesFit:
    def test_tbl_one_with_4_footnotes_does_not_overflow(self, tmp_path):
        pytest.importorskip("pptx")
        rng = np.random.default_rng(0)
        n = 300
        df = pd.DataFrame({
            "arm": rng.choice(["Placebo", "Treatment"], n),
            "age": rng.normal(60, 10, n),
            "sex": rng.choice(["F", "M"], n),
            "bmi": rng.normal(28, 5, n),
        })
        t = ps.tbl_one(df, by="arm").add_p().add_smd()
        # Sanity: there really are multiple footnotes here.
        assert len(t.footnotes) >= 4

        out = tmp_path / "t.pptx"
        t.to_pptx(out, slide_title="Table 1 — Baseline characteristics")
        bad = _shapes_in_bounds(out)
        assert not bad, (
            f"PPTX shapes overflowed slide bottom: {bad}"
        )

    def test_no_title_path_still_fits(self, tmp_path):
        pytest.importorskip("pptx")
        # Without a title the table can be taller; check the
        # footnote box still gets reserved space.
        df = pd.DataFrame({
            "arm": ["A", "B"] * 25,
            "x":   list(range(50)),
            "y":   list(range(50, 100)),
        })
        t = ps.tbl_one(df, by="arm").add_p()  # no .set_caption()
        out = tmp_path / "no_title.pptx"
        t.to_pptx(out)  # no slide_title argument
        bad = _shapes_in_bounds(out)
        assert not bad

    def test_many_footnotes_compress_table(self, tmp_path):
        pytest.importorskip("pptx")
        # Stack many user-supplied footnotes; box should still fit.
        df = pd.DataFrame({
            "arm": ["A", "B"] * 10,
            "x":   list(range(20)),
        })
        custom_footnotes = tuple(
            f"Footnote line number {i}" for i in range(8)
        )
        t = ps.tbl_one(df, by="arm").add_p().with_footnotes(
            list(custom_footnotes)
        )
        out = tmp_path / "many.pptx"
        t.to_pptx(out, slide_title="Many footnotes")
        bad = _shapes_in_bounds(out)
        assert not bad, f"overflow with {len(t.footnotes)} footnotes: {bad}"

    def test_plot_embedded_still_fits(self, tmp_path):
        pytest.importorskip("pptx")
        pytest.importorskip("matplotlib")
        pytest.importorskip("lifelines")
        rng = np.random.default_rng(0)
        df = pd.DataFrame({
            "arm": rng.choice(["A", "B"], 80),
            "time": rng.exponential(10, 80),
            "event": rng.integers(0, 2, 80),
        })
        t = ps.tbl_survival(df, time="time", event="event",
                            by="arm").with_km_plot(risk_times=[2, 5, 10])
        out = tmp_path / "km.pptx"
        t.to_pptx(out, slide_title="Survival summary with KM curve")
        bad = _shapes_in_bounds(out)
        assert not bad, f"overflow with embedded KM plot: {bad}"
